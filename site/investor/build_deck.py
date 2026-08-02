#!/usr/bin/env python3
"""Build Oppuna investor pitch deck (PPTX) — 14 core + 4 appendix slides, 16:9.

Source of truth for /investor/Oppuna-Investor-Pitch-Deck.pdf
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"

IVORY = RGBColor(0xF7, 0xF4, 0xEF)
SURFACE = RGBColor(0xFF, 0xFC, 0xF8)
SAGE = RGBColor(0x3D, 0x6B, 0x5A)
SAGE_DEEP = RGBColor(0x2F, 0x54, 0x46)
SAGE_SOFT = RGBColor(0xDC, 0xE8, 0xE2)
TEXT = RGBColor(0x1C, 0x24, 0x20)
MUTED = RGBColor(0x5A, 0x67, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC5, 0xD0, 0xCA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_DISPLAY = "Fraunces"
FONT_BODY = "Outfit"

MARGIN_L = Inches(0.65)
MARGIN_R = Inches(0.65)
CONTENT_W = Inches(12.0)


def set_run_font(run, name: str, size_pt: float, color: RGBColor, bold: bool = False):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = parse_xml(
                f'<a:{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{name}"/>'
            )
            rPr.append(el)
        else:
            el.set("typeface", name)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font=FONT_BODY,
    size=16,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, font, size, color, bold=bold)
    return box


def add_multitext(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, dict]],
):
    """lines: list of (text, style_kwargs)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, style) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = style.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(style.get("space_before", 0))
        p.space_after = Pt(style.get("space_after", 4))
        run = p.add_run()
        run.text = text
        set_run_font(
            run,
            style.get("font", FONT_BODY),
            style.get("size", 15),
            style.get("color", TEXT),
            bold=style.get("bold", False),
        )
    return box


def fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(shape, color)
    return shape


def add_round_rect(slide, left, top, width, height, color: RGBColor, adjust=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_solid(shape, color)
    try:
        shape.adjustments[0] = adjust
    except Exception:
        pass
    return shape


def add_notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    run = p.add_run()
    run.text = text


def link_shape(shape, url: str):
    shape.click_action.hyperlink.address = url


def new_slide(prs, bg=IVORY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg)
    return slide


def footer(slide, page: str, dark=False):
    color = SAGE_SOFT if dark else MUTED
    add_textbox(
        slide,
        MARGIN_L,
        Inches(7.12),
        Inches(5),
        Inches(0.28),
        "Oppuna · Confidential",
        size=11,
        color=color,
    )
    add_textbox(
        slide,
        Inches(11.4),
        Inches(7.12),
        Inches(1.3),
        Inches(0.28),
        page,
        size=11,
        color=color,
        align=PP_ALIGN.RIGHT,
    )


def title_block(slide, title: str, *, top=Inches(0.42), width=CONTENT_W, color=SAGE_DEEP):
    add_textbox(
        slide,
        MARGIN_L,
        top,
        width,
        Inches(1.05),
        title,
        font=FONT_DISPLAY,
        size=32,
        color=color,
        bold=True,
    )


def bullets(slide, left, top, width, height, items: list[str], *, size=15, color=TEXT, gap=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run_font(run, FONT_BODY, size, color)
    return box


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    icon = ASSETS / "icon.png"
    # Prefer cleaner crops when present
    splash = ASSETS / "screen_00_clean.jpg"
    if not splash.exists():
        splash = ASSETS / "screen_00.jpg"
    home = ASSETS / "screen_02_clean.jpg"
    if not home.exists():
        home = ASSETS / "screen_02.jpg"
    journal = ASSETS / "screen_01_clean.jpg"
    if not journal.exists():
        journal = ASSETS / "screen_01.jpg"
    qr = ASSETS / "play-qr.png"

    # ═══════════════════════════════════════════════════════════
    # 1 Cover
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs, IVORY)
    add_rect(s, 0, 0, Inches(5.4), SLIDE_H, SAGE_DEEP)
    if icon.exists():
        s.shapes.add_picture(str(icon), Inches(0.75), Inches(0.7), Inches(0.85), Inches(0.85))
    add_textbox(s, Inches(0.75), Inches(1.85), Inches(4.2), Inches(0.85), "Oppuna", font=FONT_DISPLAY, size=52, color=WHITE, bold=True)
    add_textbox(
        s,
        Inches(0.75),
        Inches(2.75),
        Inches(4.2),
        Inches(0.7),
        "Private AI for your thoughts.",
        font=FONT_DISPLAY,
        size=24,
        color=SAGE_SOFT,
        bold=True,
    )
    add_textbox(
        s,
        Inches(0.75),
        Inches(3.6),
        Inches(4.2),
        Inches(1.1),
        "Offline, multilingual and CBT-informed emotional wellness on Android",
        size=16,
        color=WHITE,
    )
    add_textbox(
        s,
        Inches(0.75),
        Inches(5.15),
        Inches(4.2),
        Inches(0.45),
        "Kaushik I B · Founder & Product/AI Architect",
        size=14,
        color=WHITE,
    )
    link = add_textbox(s, Inches(0.75), Inches(5.65), Inches(4.2), Inches(0.35), "https://oppuna.com", size=14, color=SAGE_SOFT)
    link_shape(link, "https://oppuna.com")
    add_textbox(
        s,
        Inches(0.75),
        Inches(6.85),
        Inches(4.2),
        Inches(0.35),
        "ADILAKSHMI INFOTECH PRIVATE LIMITED",
        size=10,
        color=SAGE_SOFT,
    )
    if home.exists():
        s.shapes.add_picture(str(home), Inches(6.55), Inches(0.55), height=Inches(6.4))
    add_notes(
        s,
        "SPEAKER (~40s)\n"
        "Open with the live product: Oppuna is a privacy-first, offline, multilingual and voice-enabled "
        "CBT-informed emotional-wellness application on Android.\n"
        "Investment thesis preview: consumer proof point for a privacy-first on-device CBT-informed wellness engine; "
        "Plus and institutional programs for early revenue; longer-term SDK for platform leverage.\n"
        "Boundary: supports everyday self-reflection—does not diagnose, treat or replace professional or emergency care.\n\n"
        "SOURCES\n"
        "- https://oppuna.com\n"
        "- Google Play: https://play.google.com/store/apps/details?id=com.oppuna.care\n"
        "- Cover screenshot: official Google Play listing creative, Aug 2026.",
    )

    # ═══════════════════════════════════════════════════════════
    # 2 Motivation & name
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Dignity and privacy are not premium features.")
    add_round_rect(s, MARGIN_L, Inches(1.65), CONTENT_W, Inches(1.55), SAGE_SOFT, adjust=0.06)
    add_textbox(
        s,
        Inches(0.95),
        Inches(1.9),
        Inches(11.4),
        Inches(1.15),
        "“I built Oppuna around a simple belief: receiving emotional-wellness support should not require an account, constant connectivity or surrendering private thoughts.”",
        font=FONT_DISPLAY,
        size=20,
        color=SAGE_DEEP,
        bold=True,
    )
    add_textbox(
        s,
        MARGIN_L,
        Inches(3.5),
        Inches(8.2),
        Inches(2.6),
        "Oppuna is a coined, Koraga-inspired name. The documented Koraga verb root “oppu” means “to be agreeable” or “to be in accord.” "
        "The complete word “Oppuna” is not presented as a literal dictionary translation.\n\n"
        "The name honours the language and cultural presence of the Koraga indigenous community of coastal Karnataka. "
        "The product is dedicated to the spirit of inclusion represented by that heritage—and designed for a broad audience.\n\n"
        "Interpretation: meeting a person where they are—with agreement, dignity and without judgement.",
        size=15,
        color=TEXT,
    )
    if icon.exists():
        s.shapes.add_picture(str(icon), Inches(10.3), Inches(3.7), Inches(2.0), Inches(2.0))
        add_textbox(s, Inches(9.9), Inches(5.85), Inches(2.8), Inches(0.5), "Living-leaf mark", size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, "02")
    add_notes(
        s,
        "SPEAKER (~40s)\n"
        "Keep cultural explanation respectful and brief. No costume imagery. Singular founder voice.\n\n"
        "SOURCES\n"
        "- D. N. Shankara Bhat, The Koraga Language — documents Koraga verb root oppu (“to be agreeable / to be in accord”): "
        "https://dnshankarabhat.net/wp-content/uploads/2014/11/Koraga_Language.pdf",
    )

    # ═══════════════════════════════════════════════════════════
    # 3 Problem
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Sensitive thoughts are still being asked to trust the cloud.")
    barriers = [
        ("1  Privacy", "Emotional content may pass through accounts, cloud servers and third-party AI systems."),
        ("2  Access", "Many products require connectivity, subscriptions or English-first interaction."),
        ("3  Trust", "People may hesitate to open up when they cannot control where their information goes."),
    ]
    for i, (h, body) in enumerate(barriers):
        left = MARGIN_L + Inches(i * 4.1)
        add_rect(s, left, Inches(1.7), Inches(3.85), Inches(2.55), SURFACE)
        add_rect(s, left, Inches(1.7), Inches(3.85), Inches(0.08), SAGE)
        add_textbox(s, left + Inches(0.25), Inches(2.0), Inches(3.35), Inches(0.5), h, font=FONT_DISPLAY, size=22, color=SAGE_DEEP, bold=True)
        add_textbox(s, left + Inches(0.25), Inches(2.65), Inches(3.35), Inches(1.3), body, size=15, color=TEXT)
    add_textbox(
        s,
        MARGIN_L,
        Inches(4.55),
        CONTENT_W,
        Inches(2.0),
        "India context (need ≠ Oppuna customers):\n"
        "• NMHS 2015–16: about 10.6% of Indian adults had a current mental disorder; treatment gaps commonly reported near 70–92%.\n"
        "• IAMAI–Kantar 2025: 958 million active internet users in India.\n"
        "• Statcounter (India mobile OS, Mar 2025 snapshot): Android about 94.8% of mobile OS traffic share.\n\n"
        "Oppuna’s wedge is general CBT-informed emotional wellness and self-reflection—not clinical treatment populations.",
        size=14,
        color=MUTED,
    )
    footer(s, "03")
    add_notes(
        s,
        "SPEAKER (~50s)\n"
        "Separate prevalence statistics from the commercial wedge. Do not imply every person with a mental disorder is a customer.\n\n"
        "SOURCES\n"
        "- NIMHANS National Mental Health Survey of India 2015–16 (current morbidity ~10.6%; treatment gap often cited 70–92%): "
        "WHO SEARO summary https://cdn.who.int/media/docs/default-source/searo/india/health-topic-pdf/summary.pdf ; "
        "PIB https://www.pib.gov.in/PressReleasePage.aspx?PRID=2188003\n"
        "- IAMAI–Kantar Internet in India Report 2025: 958M AIU — "
        "https://www.indiadigitalsummit.in/wp-content/uploads/2026/01/Internet-in-India-2025-Press-Release-Final.pdf\n"
        "- Statcounter Mobile OS India: https://gs.statcounter.com/os-market-share/mobile/india/\n"
        "- Product boundary: https://oppuna.com",
    )

    # ═══════════════════════════════════════════════════════════
    # 4 Solution + why now
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Private, structured wellness support can now run on the phone.")
    left_items = [
        "No account required",
        "Emotional content remains local",
        "On-device intelligence",
        "Offline and airplane-mode operation",
        "Voice and multilingual accessibility",
        "No advertising in vulnerable moments",
        "Local export and delete controls",
    ]
    add_textbox(s, MARGIN_L, Inches(1.55), Inches(7.2), Inches(0.35), "What Oppuna delivers", size=14, color=SAGE, bold=True)
    bullets(s, MARGIN_L, Inches(1.95), Inches(7.2), Inches(3.6), left_items, size=16, gap=9)
    add_rect(s, Inches(8.3), Inches(1.55), Inches(4.35), Inches(4.7), SAGE_DEEP)
    add_textbox(s, Inches(8.55), Inches(1.8), Inches(3.9), Inches(0.4), "Why now", font=FONT_DISPLAY, size=22, color=WHITE, bold=True)
    why = [
        "On-device AI capability is improving on mainstream Android devices",
        "Android has large reach across India",
        "AI-assisted self-help is becoming familiar",
        "Privacy concerns create a differentiated opening for local-first design",
    ]
    box = s.shapes.add_textbox(Inches(8.55), Inches(2.4), Inches(3.9), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(why):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run_font(run, FONT_BODY, 14, SAGE_SOFT)
    footer(s, "04")
    add_notes(
        s,
        "SPEAKER (~45s)\n"
        "Position as architecture opportunity, not hype. Do not invent a multiplied TAM from AIU × Android share.\n\n"
        "SOURCES\n"
        "- Product capabilities: https://oppuna.com ; https://oppuna.com/privacy/\n"
        "- India AIU and Android OS share: IAMAI–Kantar 2025; Statcounter India mobile OS (cited on slide 3 notes).",
    )

    # ═══════════════════════════════════════════════════════════
    # 5 CBT mechanism
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "CBT turns overwhelming thoughts into manageable actions.")
    stages = [
        ("1", "Recognise\nthe emotion"),
        ("2", "Capture the\nautomatic thought"),
        ("3", "Examine\nthe evidence"),
        ("4", "Build a\nbalanced reframe"),
        ("5", "Choose one\nmanageable action"),
        ("6", "Reflect and\nobserve patterns"),
    ]
    for i, (n, label) in enumerate(stages):
        left = MARGIN_L + Inches(i * 2.05)
        add_rect(s, left, Inches(1.55), Inches(1.9), Inches(1.85), SURFACE)
        add_textbox(s, left + Inches(0.1), Inches(1.65), Inches(1.7), Inches(0.4), n, font=FONT_DISPLAY, size=22, color=SAGE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.08), Inches(2.15), Inches(1.75), Inches(1.05), label, size=13, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(s, MARGIN_L, Inches(3.6), CONTENT_W, Inches(0.35), "Connected Oppuna practices", size=14, color=SAGE, bold=True)
    practices = [
        "Mood check-ins",
        "Guided thought records",
        "Thought-reframing exercises",
        "Guided journaling",
        "Behavioural / wellness goals",
        "Grounding and breathing",
        "Daily plan",
        "Progress summaries",
    ]
    for i, ptext in enumerate(practices):
        col = i % 4
        row = i // 4
        left = MARGIN_L + Inches(col * 3.1)
        top = Inches(4.05 + row * 0.7)
        add_rect(s, left, top, Inches(2.95), Inches(0.55), SAGE_SOFT)
        add_textbox(s, left + Inches(0.12), top + Inches(0.12), Inches(2.7), Inches(0.35), ptext, size=14, color=SAGE_DEEP, bold=True)
    add_textbox(
        s,
        MARGIN_L,
        Inches(5.7),
        CONTENT_W,
        Inches(0.7),
        "CBT-informed self-help for everyday wellness—not therapy or diagnosis.\n"
        "Oppuna helps people recognise thoughts and emotions, examine patterns, reframe unhelpful thinking, practise grounding and breathing, choose manageable actions and observe progress—privately on Android.",
        size=14,
        color=MUTED,
    )
    footer(s, "05")
    add_notes(
        s,
        "SPEAKER (~50s)\n"
        "This is the differentiation slide: structured CBT-informed mechanism, not a generic chatbot or meditation library.\n"
        "Repeat boundary language.\n\n"
        "SOURCES\n"
        "- Product feature set: https://oppuna.com\n"
        "- Positioning: CBT-informed emotional wellness (founder product brief, Aug 2026).",
    )

    # ═══════════════════════════════════════════════════════════
    # 6 Product journey (3 steps — only 3 genuine screenshots available)
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "One private space for the practices that build consistency.")
    steps = [
        ("1. Open a private space", splash, "No account. Begin on the device."),
        ("2. Check in and follow today’s care", home, "Mood, plan, talk, small actions."),
        ("3. Journal and keep it local", journal, "Thoughts stay on this device."),
    ]
    for i, (label, img, caption) in enumerate(steps):
        left = MARGIN_L + Inches(i * 4.15)
        add_textbox(s, left, Inches(1.5), Inches(3.9), Inches(0.4), label, size=15, color=SAGE_DEEP, bold=True)
        if img.exists():
            s.shapes.add_picture(str(img), left + Inches(0.55), Inches(2.0), height=Inches(4.2))
        add_textbox(s, left, Inches(6.35), Inches(3.9), Inches(0.4), caption, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, "06")
    add_notes(
        s,
        "SPEAKER (~40s)\n"
        "Walk the realistic journey with only verified screenshots from the Google Play listing. "
        "Full journey in words: mood check-in → structured exercise or daily plan → journal / thought record / breathe / ground / talk → small next action → return to observe patterns.\n"
        "Only three official listing screenshots were available; steps match those assets rather than fabricating screens.\n\n"
        "SOURCES\n"
        "- Screenshots: Google Play listing creatives for com.oppuna.care, Aug 2026\n"
        "- Journey capabilities: https://oppuna.com",
    )

    # ═══════════════════════════════════════════════════════════
    # 7 Architecture / privacy / safety
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Privacy and safety are product architecture—not policy promises.")
    # flow strip
    flow = [
        (Inches(0.65), "User"),
        (Inches(3.55), "Android app"),
        (Inches(6.45), "Local storage"),
        (Inches(9.35), "On-device AI +\nlocal safety"),
    ]
    for left, label in flow:
        add_rect(s, left, Inches(1.55), Inches(2.55), Inches(1.05), SURFACE)
        add_textbox(s, left + Inches(0.1), Inches(1.7), Inches(2.35), Inches(0.8), label, size=15, color=SAGE_DEEP, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for left in (Inches(3.2), Inches(6.1), Inches(9.0)):
        add_textbox(s, left, Inches(1.75), Inches(0.35), Inches(0.5), "→", font=FONT_DISPLAY, size=22, color=SAGE, bold=True)
    bullets(
        s,
        MARGIN_L,
        Inches(2.9),
        Inches(7.4),
        Inches(3.5),
        [
            "Journals, moods, thought records, voice notes and conversations remain on the device.",
            "On-device intelligence supports CBT-informed and general-wellness workflows.",
            "Local safety logic can present appropriate crisis resources; Oppuna is not an emergency service.",
            "Export and delete controls provide user control; offline operation is first-class.",
            "Optional device biometric / PIN app lock is available; SQLCipher-class database encryption is a funded improvement—not a current claim.",
        ],
        size=14,
        gap=8,
    )
    add_rect(s, Inches(8.4), Inches(2.9), Inches(4.25), Inches(3.4), SAGE_SOFT)
    add_textbox(s, Inches(8.65), Inches(3.1), Inches(3.85), Inches(0.4), "Shipped now vs funded", size=14, color=SAGE, bold=True)
    add_textbox(
        s,
        Inches(8.65),
        Inches(3.55),
        Inches(3.85),
        Inches(2.5),
        "Shipped: local storage, offline companion, CBT-informed workflows, crisis-resource presentation, export/delete, optional app lock.\n\n"
        "Funded: stronger on-device models, encryption hardening, broader device validation, professional CBT review, multilingual voice depth.",
        size=13,
        color=TEXT,
    )
    footer(s, "07")
    add_notes(
        s,
        "SPEAKER (~45s)\n"
        "Do not claim perfect safety, certifications or clinical validation. Separate shipped from funded.\n\n"
        "SOURCES\n"
        "- Privacy and architecture: https://oppuna.com/privacy/ ; docs/DATA_PROTECTION.md in product repo (SQLite not SQLCipher-encrypted today)\n"
        "- App lock implementation present in application code (expo-local-authentication / SecureStore)\n"
        "- Crisis resources on site: 112, Tele-MANAS, KIRAN via official government pages linked from oppuna.com",
    )

    # ═══════════════════════════════════════════════════════════
    # 8 Traction
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Early users are returning to a private, structured experience.")
    add_rect(s, MARGIN_L, Inches(1.7), Inches(5.7), Inches(2.4), SURFACE)
    add_textbox(s, Inches(0.95), Inches(2.05), Inches(5.2), Inches(0.9), "80–100", font=FONT_DISPLAY, size=48, color=SAGE_DEEP, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.95), Inches(3.1), Inches(5.2), Inches(0.5), "Early users", size=18, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.95), Inches(1.7), Inches(5.7), Inches(2.4), SURFACE)
    add_textbox(s, Inches(7.25), Inches(2.05), Inches(5.2), Inches(0.9), "~90%", font=FONT_DISPLAY, size=48, color=SAGE_DEEP, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        Inches(7.25),
        Inches(3.05),
        Inches(5.2),
        Inches(0.7),
        "Returned during the current\nobservation period",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        MARGIN_L,
        Inches(4.4),
        CONTENT_W,
        Inches(0.7),
        "Live on Google Play · Production website at oppuna.com\n"
        "Approximately 90% of tracked early users returned during the current observation period.",
        size=16,
        color=TEXT,
    )
    add_textbox(
        s,
        MARGIN_L,
        Inches(5.4),
        CONTENT_W,
        Inches(1.1),
        "Founder-reported early-user tracking across approximately 80–100 users. "
        "The current return measure is directional and is not yet equivalent to standard D7 or D30 cohort retention.",
        size=13,
        color=MUTED,
    )
    footer(s, "08")
    add_notes(
        s,
        "SPEAKER (~35s)\n"
        "Present as early signal, not product-market fit proof. Methodology detail is in Appendix 1.\n\n"
        "SOURCES\n"
        "- Founder-reported early-user tracking (Aug 2026)\n"
        "- https://play.google.com/store/apps/details?id=com.oppuna.care\n"
        "- https://oppuna.com",
    )

    # ═══════════════════════════════════════════════════════════
    # 9 Market wedge
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "India is the starting point for private, accessible CBT-informed wellness.")
    add_textbox(s, MARGIN_L, Inches(1.5), Inches(7.4), Inches(0.35), "Initial customer wedge", size=14, color=SAGE, bold=True)
    bullets(
        s,
        MARGIN_L,
        Inches(1.9),
        Inches(7.4),
        Inches(2.6),
        [
            "Android-first adults in India",
            "Privacy-conscious users",
            "People with inconsistent connectivity",
            "Users seeking structured self-reflection rather than clinical treatment",
            "Regional-language and voice-first users",
        ],
        size=15,
        gap=8,
    )
    add_rect(s, Inches(8.4), Inches(1.5), Inches(4.25), Inches(4.7), SURFACE)
    add_textbox(s, Inches(8.65), Inches(1.7), Inches(3.85), Inches(0.4), "Market framing", font=FONT_DISPLAY, size=18, color=SAGE_DEEP, bold=True)
    add_textbox(
        s,
        Inches(8.65),
        Inches(2.25),
        Inches(3.85),
        Inches(3.7),
        "Enabling market\n958M active internet users in India (IAMAI–Kantar 2025). Android dominates India mobile OS traffic (~94.8%, Statcounter Mar 2025).\n\n"
        "Initial segment\nPrivacy-seeking Android adults who want structured CBT-informed self-help that works offline.\n\n"
        "18-month target\n5,000 monthly active users with defined retention and early paid conversion.\n\n"
        "Expansion\nInstitutional pilots → embeddable SDK for partners.",
        size=13,
        color=TEXT,
    )
    add_textbox(
        s,
        MARGIN_L,
        Inches(5.0),
        Inches(7.4),
        Inches(1.5),
        "These figures are enabling-context statistics, not a multiplied TAM. "
        "Oppuna does not treat clinical prevalence as its addressable customer count. "
        "The near-term job is to prove retention and willingness to pay inside a focused wedge.",
        size=14,
        color=MUTED,
    )
    footer(s, "09")
    add_notes(
        s,
        "SPEAKER (~45s)\n"
        "Do not multiply AIU × Android share into a validated TAM. Keep segment qualitative + 18-month target quantitative.\n\n"
        "SOURCES\n"
        "- IAMAI–Kantar Internet in India 2025 (958M AIU)\n"
        "- Statcounter India mobile OS (Android ~94.8%, Mar 2025 snapshot)\n"
        "- NMHS for clinical context only (slide 3)\n"
        "- 18-month targets aligned with raise plan (slide 14)",
    )

    # ═══════════════════════════════════════════════════════════
    # 10 Business model
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "The core remains free; revenue comes from depth and distribution.")
    add_textbox(
        s,
        MARGIN_L,
        Inches(1.45),
        CONTENT_W,
        Inches(0.55),
        "Investment thesis: Oppuna is the consumer proof point for a privacy-first, on-device CBT-informed wellness engine. "
        "Consumer adoption validates the product; Oppuna Plus and institutional programs provide early revenue; the longer-term SDK creates platform leverage.",
        size=14,
        color=TEXT,
    )
    stages = [
        ("1", "Free core", "Trust, access and daily engagement", "Now"),
        ("2", "Oppuna Plus", "Depth, voice and personalisation", "Near-term"),
        ("3", "Institutional", "Employer, NGO, CSR pilots", "After retention"),
        ("4", "SDK / platform", "License the privacy-first engine", "Longer-term"),
    ]
    for i, (n, h, b, timing) in enumerate(stages):
        left = MARGIN_L + Inches(i * 3.15)
        bg = SAGE_DEEP if i == 0 else SURFACE
        fg = WHITE if i == 0 else SAGE_DEEP
        muted = SAGE_SOFT if i == 0 else MUTED
        add_rect(s, left, Inches(2.2), Inches(3.0), Inches(2.35), bg)
        add_textbox(s, left + Inches(0.2), Inches(2.35), Inches(2.6), Inches(0.35), f"{n}  {timing}", size=12, color=muted, bold=True)
        add_textbox(s, left + Inches(0.2), Inches(2.75), Inches(2.6), Inches(0.55), h, font=FONT_DISPLAY, size=18, color=fg, bold=True)
        add_textbox(s, left + Inches(0.2), Inches(3.45), Inches(2.6), Inches(0.8), b, size=13, color=muted)
    add_textbox(
        s,
        MARGIN_L,
        Inches(4.8),
        Inches(7.5),
        Inches(1.7),
        "Free core includes: basic CBT-informed exercises · mood check-ins · thought records · journaling · breathing and grounding · "
        "basic daily plan · crisis resources · local export and delete.\n\n"
        "Future Oppuna Plus: stronger on-device models · natural low-latency voice · advanced personalised CBT-informed journeys · "
        "deeper insights · longitudinal summaries · user-controlled encrypted backup and migration.",
        size=13,
        color=TEXT,
    )
    add_rect(s, Inches(8.4), Inches(4.8), Inches(4.25), Inches(1.7), SAGE_SOFT)
    add_textbox(s, Inches(8.65), Inches(5.0), Inches(3.85), Inches(0.35), "Pricing hypothesis", size=13, color=SAGE, bold=True)
    add_textbox(s, Inches(8.65), Inches(5.35), Inches(3.85), Inches(0.45), "₹99–₹199 / month", font=FONT_DISPLAY, size=22, color=SAGE_DEEP, bold=True)
    add_textbox(s, Inches(8.65), Inches(5.9), Inches(3.85), Inches(0.4), "Not final pricing · no ads · no data sale · no per-message fees", size=12, color=MUTED)
    footer(s, "10")
    add_notes(
        s,
        "SPEAKER (~50s)\n"
        "Sequence matters: free core first; Plus and institutional after proof; SDK is longer-term leverage—not four equal immediate businesses.\n\n"
        "SOURCES\n"
        "- Business sequence and pricing hypothesis: founder strategy brief, Aug 2026\n"
        "- Free-core capabilities: https://oppuna.com",
    )

    # ═══════════════════════════════════════════════════════════
    # 11 GTM
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Trust-led distribution begins with focused, measurable channels.")
    add_rect(s, MARGIN_L, Inches(1.55), Inches(7.5), Inches(4.6), SURFACE)
    add_textbox(s, Inches(0.95), Inches(1.75), Inches(7.0), Inches(0.4), "Phase 1 — Validation", font=FONT_DISPLAY, size=20, color=SAGE_DEEP, bold=True)
    bullets(
        s,
        Inches(0.95),
        Inches(2.3),
        Inches(7.0),
        Inches(1.6),
        [
            "Google Play and app-store optimisation",
            "Founder-led community adoption",
            "Privacy, CBT self-help and regional-language educational content",
            "User referrals",
        ],
        size=15,
        gap=8,
    )
    add_textbox(s, Inches(0.95), Inches(4.1), Inches(7.0), Inches(0.4), "Phase 2 — After retention is established", font=FONT_DISPLAY, size=18, color=SAGE_DEEP, bold=True)
    bullets(
        s,
        Inches(0.95),
        Inches(4.6),
        Inches(7.0),
        Inches(1.3),
        [
            "NGO and community partnerships",
            "Employer or CSR-sponsored pilots",
            "Selected public-wellness programs",
        ],
        size=15,
        gap=8,
    )
    add_rect(s, Inches(8.5), Inches(1.55), Inches(4.15), Inches(4.6), SAGE)
    add_textbox(s, Inches(8.8), Inches(2.1), Inches(3.6), Inches(0.4), "Early checkpoint", size=14, color=SAGE_SOFT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.8), Inches(2.7), Inches(3.6), Inches(0.9), "1,000", font=FONT_DISPLAY, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        Inches(8.8),
        Inches(3.7),
        Inches(3.6),
        Inches(1.6),
        "active users as an early validation step—not the final milestone of the ₹1 crore round.",
        size=15,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    footer(s, "11")
    add_notes(
        s,
        "SPEAKER (~40s)\n"
        "Do not invent signed partnerships. Phase 2 follows retention proof.\n\n"
        "SOURCES\n"
        "- GTM phasing and 1,000-user checkpoint: founder strategy brief, Aug 2026\n"
        "- Current channel: Google Play listing for com.oppuna.care",
    )

    # ═══════════════════════════════════════════════════════════
    # 12 Competition
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Oppuna competes through architecture and structured engagement.")
    headers = ["Dimension", "Oppuna", "Wysa", "Calm / Headspace"]
    rows = [
        ["Where content is processed", "On device", "Cloud / AWS servers", "Cloud + content CDN"],
        ["Interactive support offline", "Yes (core)", "Requires connectivity", "Downloaded content may work offline; interactive cloud features do not"],
        ["Structured CBT-informed workflows", "Yes", "Yes (care library)", "Mindfulness / sleep content focus"],
        ["Account required", "No", "Anonymous options; still server-based", "Account typical"],
        ["Voice + multilingual", "Voice roadmap; EN · HI · ES", "Chat-led; multilingual options", "Content libraries; account ecosystems"],
        ["Local data control", "Local-first", "Server-stored (transit/at-rest encryption)", "Cloud sync / accounts"],
        ["Business model", "Free → Plus → institutional → SDK", "Consumer + enterprise wellbeing", "Subscription content platforms"],
    ]
    col_w = [3.1, 2.9, 3.1, 3.1]
    x0 = 0.55
    y0 = 1.45
    x = x0
    for i, h in enumerate(headers):
        add_rect(s, Inches(x), Inches(y0), Inches(col_w[i]), Inches(0.4), SAGE_DEEP)
        add_textbox(s, Inches(x + 0.05), Inches(y0 + 0.07), Inches(col_w[i] - 0.1), Inches(0.3), h, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += col_w[i]
    for r, row in enumerate(rows):
        x = x0
        bg = SURFACE if r % 2 == 0 else SAGE_SOFT
        for i, cell in enumerate(row):
            add_rect(s, Inches(x), Inches(y0 + 0.4 + r * 0.62), Inches(col_w[i]), Inches(0.62), bg)
            add_textbox(
                s,
                Inches(x + 0.06),
                Inches(y0 + 0.45 + r * 0.62),
                Inches(col_w[i] - 0.12),
                Inches(0.52),
                cell,
                size=11,
                color=SAGE_DEEP if i == 1 else TEXT,
                bold=(i == 0 or i == 1),
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER,
            )
            x += col_w[i]
    add_textbox(
        s,
        MARGIN_L,
        Inches(6.35),
        CONTENT_W,
        Inches(0.5),
        "Emerging moat: on-device CBT-informed workflow engine · local-model optimisation · safety-aware offline workflows · voice and multilingual UX · privacy-led brand · future embeddable SDK.",
        size=12,
        color=MUTED,
    )
    footer(s, "12")
    add_notes(
        s,
        "SPEAKER (~50s)\n"
        "Correct prior overclaims: Calm/Headspace can offer offline playback of downloaded content; do not claim they have ads as a core model. "
        "Wysa offers anonymous usage but processes conversations on cloud servers (AWS) per its privacy policy.\n\n"
        "SOURCES\n"
        "- Wysa Privacy Policy (cloud/AWS storage; AI/LLM processing): https://legal.wysa.io/privacy-policy\n"
        "- Wysa FAQ (anonymous usage; connectivity/service model): https://www.wysa.com/faq\n"
        "- Calm Premium / content platform: https://support.calm.com/hc/en-us/articles/360008536834-Calm-Premium-vs-Free-Features-Content-List-Benefits\n"
        "- Headspace account/subscription model: https://www.headspace.com/\n"
        "- Calm/Headspace offline downloads commonly documented in vendor help centres (downloaded sessions may work offline)\n"
        "- Oppuna local-first claims: https://oppuna.com/privacy/",
    )

    # ═══════════════════════════════════════════════════════════
    # 13 Solo founder
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Built by a founder who takes complex systems to production.")
    add_rect(s, MARGIN_L, Inches(1.55), CONTENT_W, Inches(4.85), SURFACE)
    add_rect(s, MARGIN_L, Inches(1.55), Inches(0.16), Inches(4.85), SAGE)
    add_textbox(s, Inches(1.1), Inches(1.8), Inches(11), Inches(0.55), "Kaushik I B", font=FONT_DISPLAY, size=30, color=SAGE_DEEP, bold=True)
    add_textbox(s, Inches(1.1), Inches(2.4), Inches(11), Inches(0.4), "Founder & Product/AI Architect", size=18, color=SAGE, bold=True)
    bullets(
        s,
        Inches(1.1),
        Inches(3.0),
        Inches(11),
        Inches(2.4),
        [
            "12+ years across software architecture, cloud and AI",
            "Enterprise systems across banking, FinTech and AI platforms",
            "TOGAF certified; Microsoft Teams ecosystem experience",
            "Currently building enterprise AI and agentic systems",
            "Built Oppuna end-to-end—from architecture and on-device workflows to Google Play launch",
        ],
        size=16,
        gap=8,
    )
    add_textbox(
        s,
        Inches(1.1),
        Inches(5.5),
        Inches(11),
        Inches(0.6),
        "Kaushik combines enterprise architecture experience with the ability to build and ship privacy-sensitive AI products independently.",
        font=FONT_DISPLAY,
        size=16,
        color=SAGE_DEEP,
        bold=True,
    )
    footer(s, "13")
    add_notes(
        s,
        "SPEAKER (~40s)\n"
        "Solo founder. Do not invent additional founders, advisors or clinical staff. Emphasise capital-efficient execution and enterprise-grade systems experience applied to privacy-sensitive consumer AI.\n\n"
        "SOURCES\n"
        "- Founder credentials and narrative: founder brief, Aug 2026\n"
        "- Execution evidence: live Android app on Google Play; website https://oppuna.com",
    )

    # ═══════════════════════════════════════════════════════════
    # 14 Raise
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs, SAGE_DEEP)
    add_textbox(
        s,
        MARGIN_L,
        Inches(0.35),
        CONTENT_W,
        Inches(0.95),
        "₹1 crore takes Oppuna from early usage to validated, revenue-generating wellness infrastructure.",
        font=FONT_DISPLAY,
        size=26,
        color=WHITE,
        bold=True,
    )
    add_rect(s, MARGIN_L, Inches(1.4), Inches(5.9), Inches(1.15), RGBColor(0x3A, 0x5F, 0x52))
    add_textbox(s, Inches(0.9), Inches(1.55), Inches(5.4), Inches(0.3), "Raising", size=13, color=SAGE_SOFT)
    add_textbox(s, Inches(0.9), Inches(1.85), Inches(5.4), Inches(0.45), "₹1 crore pre-seed", font=FONT_DISPLAY, size=22, color=WHITE, bold=True)
    add_rect(s, Inches(6.85), Inches(1.4), Inches(5.8), Inches(1.15), RGBColor(0x3A, 0x5F, 0x52))
    add_textbox(s, Inches(7.1), Inches(1.55), Inches(5.4), Inches(0.3), "Runway", size=13, color=SAGE_SOFT)
    add_textbox(s, Inches(7.1), Inches(1.85), Inches(5.4), Inches(0.45), "18 months", font=FONT_DISPLAY, size=22, color=WHITE, bold=True)

    add_textbox(s, MARGIN_L, Inches(2.75), Inches(4.0), Inches(0.3), "Product & safety", size=13, color=SAGE_SOFT, bold=True)
    bullets(
        s,
        MARGIN_L,
        Inches(3.05),
        Inches(4.0),
        Inches(1.7),
        [
            "Professionally reviewed CBT journeys",
            "On-device AI + multilingual voice upgrade",
            "Local-data protection & safety testing",
            "Broader Android device validation",
        ],
        size=12,
        color=WHITE,
        gap=4,
    )
    add_textbox(s, Inches(4.85), Inches(2.75), Inches(4.0), Inches(0.3), "Engagement", size=13, color=SAGE_SOFT, bold=True)
    bullets(
        s,
        Inches(4.85),
        Inches(3.05),
        Inches(4.0),
        Inches(1.7),
        [
            "5,000 monthly active users",
            "Defined D30 retention of 25%+",
            "1,000 users completing ≥4 structured CBT-informed exercises",
        ],
        size=12,
        color=WHITE,
        gap=4,
    )
    add_textbox(s, Inches(9.0), Inches(2.75), Inches(3.7), Inches(0.3), "Commercial", size=13, color=SAGE_SOFT, bold=True)
    bullets(
        s,
        Inches(9.0),
        Inches(3.05),
        Inches(3.7),
        Inches(1.5),
        [
            "500 Oppuna Plus subscribers or ₹1 lakh MRR",
            "2–3 paid institutional pilots",
        ],
        size=12,
        color=WHITE,
        gap=4,
    )

    add_textbox(s, MARGIN_L, Inches(4.7), CONTENT_W, Inches(0.3), "Use of funds", size=13, color=SAGE_SOFT, bold=True)
    add_textbox(
        s,
        MARGIN_L,
        Inches(5.05),
        Inches(9.5),
        Inches(0.85),
        "40% Product, Android and on-device AI engineering · 15% CBT content and professional review · "
        "15% Privacy, safety, security and testing · 20% User research, community distribution and growth · "
        "10% Institutional pilots, legal and operations",
        size=13,
        color=WHITE,
    )
    add_textbox(
        s,
        MARGIN_L,
        Inches(5.9),
        Inches(9.5),
        Inches(0.55),
        "Oppuna is building a future where structured emotional-wellness support does not require surrendering privacy.",
        font=FONT_DISPLAY,
        size=14,
        color=SAGE_SOFT,
        bold=True,
    )
    web = add_textbox(s, MARGIN_L, Inches(6.55), Inches(3.6), Inches(0.35), "https://oppuna.com", size=13, color=WHITE)
    link_shape(web, "https://oppuna.com")
    mail = add_textbox(s, Inches(4.3), Inches(6.55), Inches(3.6), Inches(0.35), "support@oppuna.com", size=13, color=SAGE_SOFT)
    link_shape(mail, "mailto:support@oppuna.com")
    if qr.exists():
        pic = s.shapes.add_picture(str(qr), Inches(11.2), Inches(5.35), Inches(1.45), Inches(1.45))
        link_shape(pic, "https://play.google.com/store/apps/details?id=com.oppuna.care")
        add_textbox(s, Inches(11.05), Inches(6.85), Inches(1.75), Inches(0.25), "Google Play", size=11, color=SAGE_SOFT, align=PP_ALIGN.CENTER)
    footer(s, "14", dark=True)
    add_notes(
        s,
        "SPEAKER (~55s)\n"
        "Close on the ₹1 crore pre-seed ask and 18-month milestones. No equity % or valuation. CTA: fundraising conversation.\n\n"
        "SOURCES\n"
        "- Raise plan and milestones: founder fundraising brief, Aug 2026\n"
        "- Contact: https://oppuna.com ; mailto:support@oppuna.com\n"
        "- Play: https://play.google.com/store/apps/details?id=com.oppuna.care",
    )

    # ═══════════════════════════════════════════════════════════
    # A1 Retention methodology
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Retention methodology")
    add_textbox(s, MARGIN_L, Inches(1.4), CONTENT_W, Inches(0.3), "Appendix 1", size=13, color=SAGE, bold=True)
    bullets(
        s,
        MARGIN_L,
        Inches(1.9),
        CONTENT_W,
        Inches(4.6),
        [
            "Current claim: approximately 90% of tracked early users returned during the observation period.",
            "Population: approximately 80–100 early users.",
            "Source: founder-reported internal tracking.",
            "Limitation: not yet a frozen D1, D7 or D30 install cohort with a pre-declared return definition.",
            "Prospective framework: cohort start date · observation window · qualifying return event · exclusions · D1/D7/D30 · monthly active users.",
            "Purpose of the upgrade: make retention comparable for diligence and milestone tracking against the 25%+ D30 target.",
        ],
        size=16,
        gap=10,
    )
    footer(s, "A1")
    add_notes(
        s,
        "SPEAKER (optional)\n"
        "Investor diligence detail for the early return measure.\n\n"
        "SOURCES\n"
        "- Founder-reported early-user tracking, Aug 2026\n"
        "- Future source of truth: Google Play Console retention / cohort reports",
    )

    # ═══════════════════════════════════════════════════════════
    # A2 Architecture current vs funded
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Current architecture and funded improvements")
    add_textbox(s, MARGIN_L, Inches(1.4), CONTENT_W, Inches(0.3), "Appendix 2", size=13, color=SAGE, bold=True)
    add_rect(s, MARGIN_L, Inches(1.85), Inches(5.85), Inches(4.7), SURFACE)
    add_textbox(s, Inches(0.9), Inches(2.05), Inches(5.3), Inches(0.4), "Current product", font=FONT_DISPLAY, size=20, color=SAGE_DEEP, bold=True)
    bullets(
        s,
        Inches(0.9),
        Inches(2.6),
        Inches(5.3),
        Inches(3.7),
        [
            "Android application",
            "Local data storage",
            "On-device / offline companion",
            "CBT-informed wellness workflows",
            "Mood, journal and thought records",
            "Local crisis-resource presentation",
            "Export / delete controls",
            "Offline operation",
            "Optional biometric / PIN app lock",
        ],
        size=14,
        gap=6,
    )
    add_rect(s, Inches(6.85), Inches(1.85), Inches(5.8), Inches(4.7), SAGE_SOFT)
    add_textbox(s, Inches(7.1), Inches(2.05), Inches(5.3), Inches(0.4), "Funded improvements", font=FONT_DISPLAY, size=20, color=SAGE_DEEP, bold=True)
    bullets(
        s,
        Inches(7.1),
        Inches(2.6),
        Inches(5.3),
        Inches(3.7),
        [
            "Stronger on-device models",
            "Encryption hardening",
            "Biometric / PIN protection hardening",
            "Broader device validation",
            "Professional CBT-content review",
            "Expanded safety testing",
            "Multilingual voice improvement",
            "User-controlled backup and migration",
        ],
        size=14,
        gap=6,
    )
    footer(s, "A2")
    add_notes(
        s,
        "SPEAKER (optional)\n"
        "Clear separation of shipped vs funded. Database encryption is funded, not shipped.\n\n"
        "SOURCES\n"
        "- https://oppuna.com/privacy/\n"
        "- Product data-protection notes: SQLite currently not SQLCipher-encrypted; app lock exists via device authentication APIs",
    )

    # ═══════════════════════════════════════════════════════════
    # A3 Roadmap
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Product and revenue roadmap")
    add_textbox(s, MARGIN_L, Inches(1.4), CONTENT_W, Inches(0.3), "Appendix 3", size=13, color=SAGE, bold=True)
    cols = [
        ("Free core", "Shipped now", ["CBT-informed basics", "Mood + thought records", "Journaling", "Breath / ground / sleep", "Daily plan", "Crisis resources", "Export / delete"]),
        ("Oppuna Plus", "In development / hypothesis", ["Stronger on-device models", "Natural voice", "Personalised journeys", "Deeper insights", "Longitudinal summaries", "Encrypted backup*", "Device migration*"]),
        ("Institutional", "Future", ["Sponsored access", "Employer pilots", "NGO / CSR programs", "Public-wellness pilots", "No sale of user data"]),
        ("SDK / platform", "Longer-term", ["Privacy-first engine", "White-label path", "Partner integrations", "Local-model packaging", "Safety workflows"]),
    ]
    for i, (h, status, items) in enumerate(cols):
        left = MARGIN_L + Inches(i * 3.15)
        add_rect(s, left, Inches(1.8), Inches(3.0), Inches(4.7), SURFACE if i % 2 == 0 else SAGE_SOFT)
        add_textbox(s, left + Inches(0.15), Inches(1.95), Inches(2.7), Inches(0.4), h, size=15, color=SAGE_DEEP, bold=True)
        add_textbox(s, left + Inches(0.15), Inches(2.35), Inches(2.7), Inches(0.35), status, size=11, color=SAGE, bold=True)
        bullets(s, left + Inches(0.1), Inches(2.8), Inches(2.8), Inches(3.4), items, size=12, gap=5)
    add_textbox(
        s,
        MARGIN_L,
        Inches(6.65),
        CONTENT_W,
        Inches(0.3),
        "* Future / hypothesis — not shipped. Roadmap items are not existing features.",
        size=12,
        color=MUTED,
    )
    footer(s, "A3")
    add_notes(
        s,
        "SPEAKER (optional)\n"
        "Status labels matter for diligence: shipped vs in development vs future vs hypothesis.\n\n"
        "SOURCES\n"
        "- Shipped capabilities: https://oppuna.com\n"
        "- Roadmap sequencing: founder strategy brief, Aug 2026",
    )

    # ═══════════════════════════════════════════════════════════
    # A4 Sources
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    title_block(s, "Sources")
    add_textbox(s, MARGIN_L, Inches(1.4), CONTENT_W, Inches(0.3), "Appendix 4", size=13, color=SAGE, bold=True)
    sources = [
        "Koraga verb root “oppu”: D. N. Shankara Bhat, The Koraga Language — https://dnshankarabhat.net/wp-content/uploads/2014/11/Koraga_Language.pdf",
        "India mental-health context: NIMHANS National Mental Health Survey 2015–16 — WHO SEARO summary; PIB releases citing NMHS.",
        "India internet users: IAMAI–Kantar Internet in India Report 2025 (958M active internet users).",
        "Android share (India mobile OS traffic): Statcounter Global Stats — https://gs.statcounter.com/os-market-share/mobile/india/",
        "Oppuna product, privacy and positioning: https://oppuna.com · https://oppuna.com/privacy/",
        "Google Play listing: https://play.google.com/store/apps/details?id=com.oppuna.care",
        "Wysa cloud/server processing and privacy: https://legal.wysa.io/privacy-policy · https://www.wysa.com/faq",
        "Calm / Headspace: subscription content platforms; vendor help centres document offline playback of downloaded content.",
        "Traction (80–100 users; ~90% return during observation): founder-reported early-user tracking, Aug 2026.",
        "Raise plan (₹1 crore / 18 months / milestones): founder fundraising brief, Aug 2026.",
    ]
    bullets(s, MARGIN_L, Inches(1.8), CONTENT_W, Inches(4.9), sources, size=12, gap=6)
    footer(s, "A4")
    add_notes(
        s,
        "SPEAKER (optional)\n"
        "Point investors to primary sources. Prefer official pages over secondary blogs.\n\n"
        "Additional references\n"
        "- WHO SEARO NMHS summary PDF\n"
        "- IAMAI–Kantar press PDF (India Digital Summit, Jan 2026)\n"
        "- Calm Premium help article; Headspace site for account/subscription model",
    )

    out = ROOT / "Oppuna-Investor-Pitch-Deck.pptx"
    prs.save(str(out))
    print(f"Wrote {out} ({len(prs.slides)} slides)")
    return out


if __name__ == "__main__":
    build()
